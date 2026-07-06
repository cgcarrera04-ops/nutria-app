#!/usr/bin/env python3
"""Extract text from NOTA1 PPTX and all Semana PDFs for reference."""
import os, sys, io

# Force UTF-8 output
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ─── PPTX extraction ───
from pptx import Presentation
pptx_path = r"D:\Descargas 2.0\PROYECTO\NOTA_1_Semana_4_Grupo_X_actualizada.pptx"
prs = Presentation(pptx_path)
print("="*80)
print("NOTA_1_Semana_4_Grupo_X_actualizada.pptx")
print("="*80)
for i, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip().replace('\n',' ') for cell in row.cells]
                print(" | ".join(cells))

# ─── PDF extraction ───
from PyPDF2 import PdfReader

pdf_files = [
    r"D:\Descargas 2.0\PROYECTO\Semana 3 - Business Case y Gestion de Interesados (1).pdf",
    r"D:\Descargas 2.0\PROYECTO\Semana 5 - Planificacion del Cronograma (1).pdf",
    r"D:\Descargas 2.0\PROYECTO\Semana 7 - Incertidumbre y Riesgos 2026-1.pdf",
]

# Also try files with special chars
import glob
for pattern in ["Semana 4*pdf", "Semana 6*pdf"]:
    found = glob.glob(os.path.join(r"D:\Descargas 2.0\PROYECTO", pattern))
    pdf_files.extend(found)

for pdf_path in pdf_files:
    if not os.path.exists(pdf_path):
        continue
    print("\n" + "="*80)
    print(os.path.basename(pdf_path))
    print("="*80)
    try:
        reader = PdfReader(pdf_path)
        for j, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                lines = text.strip().split("\n")
                print(f"\n--- PAGE {j+1} ---")
                for line in lines[:50]:
                    print(line)
                if len(lines) > 50:
                    print(f"... ({len(lines)-50} more lines)")
    except Exception as e:
        print(f"Error: {e}")
